import streamlit as st
from pptx import Presentation
import requests
from io import BytesIO

#
def download_presentation(url):
    response = requests.get(url)
    response.raise_for_status()  
    return BytesIO(response.content)

def load_presentation(file_path):
    return Presentation(file_path)

def display_presentation(presentation):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                st.write(shape.text)
            if shape.shape_type == 13:  
                image = shape.image
                with st.container():
                    st.image(image.blob)
            if shape.has_table:
                table = shape.table
                st.write("Table:")
                for row in table.rows:
                    st.write([cell.text for cell in row.cells])
            if shape.has_chart:
                chart = shape.chart
                
                chart_data = chart.chart_data
                st.write("Chart data:")
                for series in chart_data.series:
                    st.write(f"Series {series.name}:")
                    for category, value in zip(chart_data.categories, series.values):
                        st.write(f"{category}: {value}")

# Apply custom CSS to change the background color and make font bold
st.markdown(
    """
    <style>
    .main {
        background-color: #656D4E;
        color: white;  /* Optional: Set text color to white for better contrast */
    }
    .main * {
        font-weight: bold;
    }
    </style>
    """,
    unsafe_allow_html=True
)


USERNAME = "Gaza"
PASSWORD = 123  

# Create a login function
def login():
    st.title("Login")  
    st.session_state['logged_in'] = False
    username = st.text_input("Username")
    password = st.text_input("Password", type="password")
    if st.button("Login"):
        try:
            password = int(password)  
            if username == USERNAME and password == PASSWORD:
                st.session_state['logged_in'] = True
            else:
                st.error("Invalid username or password")
        except ValueError:
            st.error("Password must be a number")

# Main function to run the app
def main():
    if 'logged_in' not in st.session_state:
        st.session_state['logged_in'] = False

    if not st.session_state['logged_in']:
        login()
    else:
        file_url = "https://github.com/mahmoudsayedashour/Social-Media-Effects/blob/main/BI%20Service%20final%20%20001.pptx"  
        file_path = download_presentation(file_url)
        presentation = load_presentation(file_path)
        display_presentation(presentation)

if __name__ == "__main__":
    main()
